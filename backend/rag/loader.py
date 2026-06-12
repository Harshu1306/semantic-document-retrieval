"""
Document loaders for PDF, DOCX, TXT, PPTX.
Returns list of LangChain Document objects with metadata.
"""

import logging
import os
from pathlib import Path
from typing import List

from langchain_core.documents import Document

logger = logging.getLogger(__name__)


def load_pdf(file_path: str) -> List[Document]:
    from langchain_community.document_loaders import PyPDFLoader

    loader = PyPDFLoader(file_path)
    return loader.load()


def load_docx(file_path: str) -> List[Document]:
    """
    Load DOCX content including paragraphs and tables.
    """

    from docx import Document as DocxDocument

    doc = DocxDocument(file_path)

    content = []

    # Extract normal paragraphs
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            content.append(text)

    # Extract tables
    for table in doc.tables:
        for row in table.rows:
            row_data = []

            for cell in row.cells:
                cell_text = cell.text.strip()

                if cell_text:
                    row_data.append(cell_text)

            if row_data:
                content.append(" | ".join(row_data))

    full_text = "\n".join(content)

    logger.info(
        "DOCX extracted %d characters from '%s'",
        len(full_text),
        os.path.basename(file_path),
    )

    return [
        Document(
            page_content=full_text,
            metadata={
                "page": 1,
                "source": file_path,
            },
        )
    ]


def load_txt(file_path: str) -> List[Document]:
    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        text = f.read()

    return [
        Document(
            page_content=text,
            metadata={
                "page": 1,
                "source": file_path,
            },
        )
    ]


def load_pptx(file_path: str) -> List[Document]:
    from pptx import Presentation

    prs = Presentation(file_path)

    docs = []

    for slide_num, slide in enumerate(prs.slides, start=1):

        texts = []

        for shape in slide.shapes:

            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text)

        if texts:
            docs.append(
                Document(
                    page_content="\n".join(texts),
                    metadata={
                        "page": slide_num,
                        "source": file_path,
                    },
                )
            )

    return docs


LOADER_MAP = {
    ".pdf": load_pdf,
    ".docx": load_docx,
    ".doc": load_docx,
    ".txt": load_txt,
    ".md": load_txt,
    ".pptx": load_pptx,
    ".ppt": load_pptx,
}


def load_document(file_path: str) -> List[Document]:
    """
    Dispatch to the appropriate loader.
    """

    ext = Path(file_path).suffix.lower()

    loader_fn = LOADER_MAP.get(ext)

    if loader_fn is None:
        raise ValueError(
            f"Unsupported file type: {ext}. Supported: {list(LOADER_MAP.keys())}"
        )

    docs = loader_fn(file_path)

    logger.info(
        "Loaded %d document(s) from '%s'",
        len(docs),
        os.path.basename(file_path),
    )

    return docs